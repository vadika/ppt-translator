import argparse
from pptx import Presentation
from pptx.shapes.placeholder import PlaceholderGraphicFrame
import os
import anthropic

def translate_text(text, target_lang, shape_type=""):
    if not text.strip():  # Skip empty or whitespace-only text
        return text
        
    print(f"\nTranslating {shape_type}:")
    print(f"Original: '{text[:50]}{'...' if len(text) > 50 else ''}'")
    
    # Get API key from environment
    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        raise ValueError("ANTHROPIC_API_KEY environment variable not set")
    
    # Map language codes to full names
    lang_map = {
        'ru': 'Russian',
        'fi': 'Finnish',
        'et': 'Estonian',
        'sv': 'Swedish',
        'en': 'English'
    }
    
    client = anthropic.Anthropic(api_key=api_key)
    try:
        prompt = f"Translate the following text to {lang_map[target_lang]}. Only return the translation, no explanations:\n\n{text}"
        message = client.messages.create(
            model="claude-3-opus-20240229",
            max_tokens=1000,
            temperature=0,
            messages=[{"role": "user", "content": prompt}]
        )
        # Extract the text content and ensure it's a string
        translated = message.content[0].text if isinstance(message.content, list) else str(message.content)
        print(f"Translated: '{translated[:50]}{'...' if len(translated) > 50 else ''}'")
        return translated
    except Exception as e:
        print(f"Translation error: {e}")
        return text

def translate_table(table, target_lang):
    """Translate text in table cells"""
    for row in table.rows:
        for cell in row.cells:
            if cell.text.strip():
                cell.text = translate_text(cell.text, target_lang, "Table cell")

def translate_smartart(shape, target_lang):
    """Recursively translate text in SmartArt graphics"""
    # Handle text in the shape itself
    if hasattr(shape, 'text'):
        shape_type = "SmartArt text"
        shape.text = translate_text(shape.text, target_lang, shape_type)
    
    # Handle text in text frames
    if hasattr(shape, 'text_frame'):
        shape_type = "SmartArt text frame"
        shape.text_frame.text = translate_text(shape.text_frame.text, target_lang, shape_type)
        
        # Process paragraphs in text frame
        for paragraph in shape.text_frame.paragraphs:
            paragraph.text = translate_text(paragraph.text, target_lang, "SmartArt paragraph")
    
    # Process child shapes if they exist
    if hasattr(shape, 'shapes'):
        for child_shape in shape.shapes:
            translate_smartart(child_shape, target_lang)
            
    # Handle SmartArt data
    if hasattr(shape, 'placeholders'):
        for placeholder in shape.placeholders:
            if hasattr(placeholder, 'text'):
                placeholder.text = translate_text(placeholder.text, target_lang, "SmartArt placeholder")

def debug_print(message, verbose):
    if verbose:
        print(message)

def translate_presentation(input_file, target_lang, verbose=False):
    # Open presentation
    print(f"\nOpening presentation: {input_file}")
    prs = Presentation(input_file)
    
    # Translate each slide's text content
    total_slides = len(prs.slides)
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\nProcessing slide {slide_num}/{total_slides}")
        for shape in slide.shapes:
            debug_print(f"Debug: Shape type = {shape.shape_type}, Shape name = {shape.name}, Shape class = {type(shape).__name__}", verbose)
            
            # Handle different types of SmartArt and Placeholders
            if shape.shape_type in [6, 7]:  # GROUP, DIAGRAM
                translate_smartart(shape, target_lang)
            elif shape.shape_type == 14 or shape.shape_type is None:  # PLACEHOLDER or shape with None type
                if hasattr(shape, "text"):
                    shape_type = "Placeholder text"
                    shape.text = translate_text(shape.text, target_lang, shape_type)
                # Recursively process shapes inside placeholder
                if hasattr(shape, 'shapes'):
                    for subshape in shape.shapes:
                        debug_print(f"Debug: Placeholder subshape type = {subshape.shape_type}, Shape name = {subshape.name}, Shape class = {type(subshape).__name__}", verbose)
                        if subshape.shape_type in [6, 7]:
                            translate_smartart(subshape, target_lang)
                        elif subshape.shape_type == 19:
                            translate_table(subshape.table, target_lang)
                        elif hasattr(subshape, "text"):
                            shape_type = type(subshape).__name__
                            subshape.text = translate_text(subshape.text, target_lang, shape_type)
            elif shape.shape_type == 19:  # TABLE
                translate_table(shape.table, target_lang)
            elif hasattr(shape, "text"):
                shape_type = type(shape).__name__
                shape.text = translate_text(shape.text, target_lang, shape_type)
            
            # Check for nested shapes in other shape types
            if hasattr(shape, 'shapes') and shape.shape_type != 14:  # Skip placeholders as they're handled above
                for subshape in shape.shapes:
                    debug_print(f"Debug: Subshape type = {subshape.shape_type}, Shape name = {subshape.name}, Shape class = {type(subshape).__name__}", verbose)
                    if subshape.shape_type in [6, 7]:
                        translate_smartart(subshape, target_lang)
    
    # Generate output filename
    file_name, file_ext = os.path.splitext(input_file)
    output_file = f"{file_name}-{target_lang}{file_ext}"
    
    # Save translated presentation
    prs.save(output_file)
    return output_file

def main():
    parser = argparse.ArgumentParser(description='Translate PowerPoint content to specified language')
    parser.add_argument('input_file', help='Input PowerPoint file')
    parser.add_argument('language', choices=['ru', 'fi', 'et', 'sv', 'en'], 
                        help='Target language (ru=Russian, fi=Finnish, et=Estonian, sv=Swedish, en=English)')
    parser.add_argument('-v', '--verbose', action='store_true',
                        help='Enable verbose debug output')
    
    args = parser.parse_args()
    
    print(f"\nStarting translation to {args.language.upper()}...")
    output_file = translate_presentation(args.input_file, args.language, args.verbose)
    print(f"\nTranslation complete!")
    print(f"Translated presentation saved as: {output_file}")

if __name__ == "__main__":
    main()
